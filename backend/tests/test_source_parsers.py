from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from app.sources.parsers.docx import parse_docx
from app.sources.parsers.office import (
    OfficeConversionUnavailable,
    convert_legacy_office,
    render_pptx_previews,
)
from app.sources.parsers.pdf import parse_pdf
from app.sources.parsers.pptx import parse_pptx
from app.sources.parsers.text import parse_text


FIXTURES = Path(__file__).parent / "fixtures" / "sources"


def test_pdf_parser_preserves_page_order_and_renders_previews(tmp_path: Path):
    parsed = parse_pdf(FIXTURES / "sample.pdf", tmp_path)

    assert parsed.page_count == 2
    assert [block.page_number for block in parsed.blocks] == [1, 2]
    assert [block.text for block in parsed.blocks] == ["First page definition", "Second page formula"]
    assert parsed.render_paths == ("pages/page-0001.png", "pages/page-0002.png")
    assert all((tmp_path / relative).is_file() for relative in parsed.render_paths)


def test_docx_parser_preserves_headings_tables_and_inline_images(tmp_path: Path):
    parsed = parse_docx(FIXTURES / "sample.docx", tmp_path)
    kinds = [block.kind for block in parsed.blocks]

    assert parsed.page_count is None
    assert kinds == ["heading", "paragraph", "table", "image"]
    assert parsed.blocks[0].text == "Semiconductor Fundamentals"
    assert parsed.blocks[1].heading_path == ("Semiconductor Fundamentals",)
    assert "Band gap | Energy range" in parsed.blocks[2].text
    assert parsed.blocks[3].asset_path == "images/image-0001.png"
    assert (tmp_path / "images" / "image-0001.png").is_file()


def test_pptx_parser_preserves_slide_order_and_speaker_notes(tmp_path: Path):
    parsed = parse_pptx(FIXTURES / "sample.pptx", tmp_path, renderer=lambda _source, _output: ())
    texts = [block.text for block in parsed.blocks]

    assert parsed.page_count == 2
    assert texts.index("Slide One") < texts.index("Speaker note for slide one") < texts.index("Slide Two")
    assert parsed.blocks[0].locator.startswith("slide:1")
    assert any("PPTX 预览不可用" in warning for warning in parsed.warnings)
    assert parsed.render_paths == ()


def test_pptx_parser_indexes_tables_and_image_only_slides(tmp_path: Path):
    source = tmp_path / "coverage.pptx"
    image = tmp_path / "pixel.png"
    Image.new("RGB", (2, 2), "white").save(image)
    presentation = Presentation()
    table_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = table_slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "参数"
    table.cell(0, 1).text = "数值"
    image_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_slide.shapes.add_picture(str(image), Inches(1), Inches(1))
    presentation.save(source)

    parsed = parse_pptx(source, tmp_path / "assets", renderer=lambda _source, _output: ())

    assert parsed.page_count == 2
    assert parsed.blocks[0].kind == "table"
    assert parsed.blocks[0].text == "参数 | 数值"
    assert parsed.blocks[1].kind == "image"
    assert parsed.blocks[1].page_number == 2


def test_powerpoint_renderer_collects_localized_export_names(tmp_path: Path):
    def localized_exporter(_source: Path, output_dir: Path) -> None:
        (output_dir / "幻灯片1.PNG").write_bytes(b"slide-one")
        (output_dir / "幻灯片2.PNG").write_bytes(b"slide-two")

    paths = render_pptx_previews(
        FIXTURES / "sample.pptx",
        tmp_path,
        exporter=localized_exporter,
    )

    assert paths == ("slides/slide-0001.png", "slides/slide-0002.png")
    assert (tmp_path / paths[0]).read_bytes() == b"slide-one"
    assert (tmp_path / paths[1]).read_bytes() == b"slide-two"


def test_text_and_markdown_parsers_decode_and_preserve_structure(tmp_path: Path):
    gb18030_path = tmp_path / "gb18030.txt"
    gb18030_path.write_bytes("定义：带隙。".encode("gb18030"))

    plain = parse_text(gb18030_path, tmp_path / "plain")
    markdown = parse_text(FIXTURES / "sample.md", tmp_path / "markdown")

    assert plain.blocks[0].text == "定义：带隙。"
    assert markdown.blocks[0].kind == "heading"
    assert markdown.blocks[0].text == "PN 结复习"
    assert [block.kind for block in markdown.blocks[1:]] == ["list", "list"]


@pytest.mark.parametrize("extension", [".doc", ".ppt"])
def test_legacy_office_without_com_has_actionable_error(tmp_path: Path, extension: str):
    source = tmp_path / f"legacy{extension}"
    source.write_bytes(b"legacy-office-placeholder")

    with pytest.raises(OfficeConversionUnavailable) as caught:
        convert_legacy_office(source, tmp_path / "converted", office_available=lambda: False)

    assert caught.value.code == "OFFICE_CONVERSION_UNAVAILABLE"
    assert caught.value.status_code == 422
    assert caught.value.action == "convert_to_modern_format"
    assert "另存为 .docx/.pptx" in caught.value.message
