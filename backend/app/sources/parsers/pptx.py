from collections.abc import Callable
from pathlib import Path

from pptx import Presentation

from app.sources.parsers.contracts import ParsedBlock, ParsedSource
from app.sources.parsers.office import render_pptx_previews


SlideRenderer = Callable[[Path, Path], tuple[str, ...]]


def parse_pptx(
    source: Path,
    output_dir: Path,
    *,
    renderer: SlideRenderer = render_pptx_previews,
) -> ParsedSource:
    output_dir.mkdir(parents=True, exist_ok=True)
    presentation = Presentation(source)
    blocks: list[ParsedBlock] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_block_start = len(blocks)
        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""
        if title:
            blocks.append(
                ParsedBlock(
                    locator=f"slide:{slide_number}:title",
                    kind="heading",
                    text=title,
                    page_number=slide_number,
                    heading_path=(title,),
                )
            )
        heading_path = (title,) if title else ()
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            kind = "paragraph"
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
            elif getattr(shape, "has_table", False):
                rows = [
                    " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    for row in shape.table.rows
                ]
                text = "\n".join(row for row in rows if row)
                kind = "table"
            else:
                continue
            if not text:
                continue
            blocks.append(
                ParsedBlock(
                    locator=f"slide:{slide_number}:shape:{shape.shape_id}",
                    kind=kind,
                    text=text,
                    page_number=slide_number,
                    heading_path=heading_path,
                )
            )
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            blocks.append(
                ParsedBlock(
                    locator=f"slide:{slide_number}:notes",
                    kind="paragraph",
                    text=notes,
                    page_number=slide_number,
                    heading_path=heading_path,
                )
            )
        if len(blocks) == slide_block_start:
            blocks.append(
                ParsedBlock(
                    locator=f"slide:{slide_number}:image",
                    kind="image",
                    text="本页无可提取文本，请结合页面预览分析。",
                    page_number=slide_number,
                    heading_path=(),
                )
            )

    warnings: tuple[str, ...] = ()
    try:
        render_paths = renderer(source, output_dir)
    except (OSError, RuntimeError):
        render_paths = ()
    if not render_paths:
        warnings = ("PPTX 预览不可用；已保留幻灯片文本和讲者备注。",)
    return ParsedSource(len(presentation.slides), tuple(blocks), warnings, render_paths)
